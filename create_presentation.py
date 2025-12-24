#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Breadtopia 프로젝트 발표용 PPTX 템플릿 생성기
빵집 테마: 따뜻한 색감, 귀여운 디자인
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 빵토피아 테마 컬러 팔레트
COLORS = {
    'primary': RGBColor(255, 216, 155),      # #FFD89B - 부드러운 빵색
    'secondary': RGBColor(255, 168, 87),     # #FFA857 - 따뜻한 오렌지
    'accent': RGBColor(245, 187, 133),       # #F5BB85 - 크루아상색
    'dark': RGBColor(139, 90, 43),           # #8B5A2B - 갈색 (텍스트)
    'light': RGBColor(255, 248, 238),        # #FFF8EE - 연한 크림색
    'white': RGBColor(255, 255, 255),        # #FFFFFF - 흰색
    'text_dark': RGBColor(51, 51, 51),       # #333333 - 진한 회색
    'text_light': RGBColor(102, 102, 102),   # #666666 - 회색
}

def create_presentation():
    """빵토피아 발표용 PPTX 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 타이틀 슬라이드
    create_title_slide(prs)

    # 2. 목차 슬라이드
    create_table_of_contents_slide(prs)

    # 3. 섹션 구분 슬라이드 (서비스 소개)
    create_section_slide(prs, "1. 서비스 소개")

    # 4. 일반 콘텐츠 슬라이드 (2열 레이아웃)
    create_two_column_slide(prs)

    # 5. AI 기능 소개 슬라이드
    create_ai_feature_slide(prs)

    # 6. 코드 예시 슬라이드
    create_code_slide(prs)

    # 7. 데이터베이스 구조 슬라이드
    create_database_slide(prs)

    # 8. 종료 슬라이드
    create_ending_slide(prs)

    # 저장
    output_path = "c:\\Users\\SSAFY\\Desktop\\1st_PJT_recordings\\Breadtopia_Presentation_Template.pptx"
    prs.save(output_path)
    print(f"PPTX file created: {output_path}")
    return output_path


def add_decorative_shapes(slide):
    """슬라이드에 장식 도형 추가 (귀여운 빵 테마)"""
    # 좌측 상단 작은 원들 (빵 부스러기 느낌)
    for i, (x, y, size) in enumerate([
        (0.2, 0.2, 0.15),
        (0.5, 0.3, 0.1),
        (9.5, 0.25, 0.12),
    ]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['accent'] if i % 2 == 0 else COLORS['primary']
        shape.line.fill.background()

    # 하단 웨이브 장식
    wave = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(6.8),
        Inches(10), Inches(0.7)
    )
    wave.fill.solid()
    wave.fill.fore_color.rgb = COLORS['light']
    wave.line.fill.background()


def create_title_slide(prs):
    """타이틀 슬라이드 - 표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 배경 그라데이션 효과 (큰 사각형)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['light']
    bg_shape.line.fill.background()

    # 상단 장식 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.5)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()

    # 프로젝트명 (큰 타이틀)
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🍞 Breadtopia (빵토피아)"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.9),
        Inches(8), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "빵지순례의 새로운 시작, 당신만의 빵 여행을 시작하세요"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLORS['text_light']
    subtitle_para.alignment = PP_ALIGN.CENTER

    # AI 기능 강조
    ai_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(5),
        Inches(5), Inches(0.6)
    )
    ai_frame = ai_box.text_frame
    ai_frame.text = "🤖 AI 기반 3대 서비스 탑재"
    ai_para = ai_frame.paragraphs[0]
    ai_para.font.size = Pt(20)
    ai_para.font.color.rgb = COLORS['secondary']
    ai_para.font.bold = True
    ai_para.alignment = PP_ALIGN.CENTER

    # 하단 정보 (팀원, 날짜)
    footer_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5),
        Inches(8), Inches(0.6)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "SSAFY 12기 최종 프로젝트 | 2025.12.24"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(16)
    footer_para.font.color.rgb = COLORS['text_light']
    footer_para.alignment = PP_ALIGN.CENTER

    # 장식 도형
    add_decorative_shapes(slide)


def create_table_of_contents_slide(prs):
    """목차 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 헤더 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "📌 목차"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 목차 내용 영역 (여기에 실제 목차를 추가하세요)
    content_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.5),
        Inches(7), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # 목차 항목들
    toc_items = [
        "1. 서비스 소개",
        "2. 개발 배경",
        "3. 기술 스택",
        "4. 시스템 아키텍처 & ERD",
        "5. 핵심 기능 소개 (AI 3대 기능)",
        "6. 필수 요구사항 충족",
        "7. AI 기능 및 추천 알고리즘 상세",
        "8. 게이미피케이션 시스템",
        "9. 주요 코드 설명",
        "10. 프로젝트 회고",
    ]

    for i, item in enumerate(toc_items):
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(12)
        p.level = 0

    # 장식
    add_decorative_shapes(slide)


def create_section_slide(prs, section_title):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 그라데이션 (큰 원형)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()

    # 중앙 원형 장식
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3), Inches(2),
        Inches(4), Inches(4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['primary']
    circle.line.width = Pt(3)
    circle.line.color.rgb = COLORS['white']

    # 섹션 제목
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER


def create_two_column_slide(prs):
    """2열 레이아웃 슬라이드 (내용 + 이미지/설명)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['light']
    header.line.fill.background()

    # 제목 영역 (사용자가 채울 부분)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "[제목을 입력하세요]"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 좌측 컬럼 (내용)
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(4.5), Inches(5.8)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.text = "[여기에 내용을 입력하세요]\n\n• 항목 1\n• 항목 2\n• 항목 3"
    for para in left_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = COLORS['text_dark']
        para.space_before = Pt(6)

    # 우측 컬럼 배경 (연한 색)
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.2),
        Inches(4.3), Inches(5.8)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = COLORS['light']
    right_bg.line.fill.background()

    # 우측 컬럼 (이미지 또는 추가 설명)
    right_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(1.5),
        Inches(3.8), Inches(5.2)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.text = "[이미지 또는 \n추가 설명 영역]\n\n💡 핵심 포인트를\n이곳에 작성하세요"
    right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in right_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = COLORS['text_light']
        para.alignment = PP_ALIGN.CENTER

    # 하단 장식
    add_decorative_shapes(slide)


def create_ai_feature_slide(prs):
    """AI 기능 소개 슬라이드 (3개 카드 레이아웃)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🤖 AI 기반 3대 핵심 기능"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    # 3개 카드 (가로 배치)
    cards = [
        {
            'title': '1. 오늘의 빵집 추천',
            'icon': '🍞',
            'desc': '취향 키워드 기반\n하루 1개 빵집 추천',
            'x': 0.5
        },
        {
            'title': '2. 웰시코기 챗봇',
            'icon': '🐶',
            'desc': '우측 하단 플로팅\n실시간 문의 응대',
            'x': 3.5
        },
        {
            'title': '3. 리뷰 요약',
            'icon': '📝',
            'desc': 'AI가 수십 개 리뷰를\n3-5줄로 요약',
            'x': 6.5
        }
    ]

    for card in cards:
        # 카드 배경
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(card['x']), Inches(1.5),
            Inches(2.8), Inches(4.5)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLORS['light']
        card_bg.line.width = Pt(2)
        card_bg.line.color.rgb = COLORS['primary']

        # 아이콘
        icon_box = slide.shapes.add_textbox(
            Inches(card['x'] + 0.3), Inches(2),
            Inches(2.2), Inches(0.8)
        )
        icon_frame = icon_box.text_frame
        icon_frame.text = card['icon']
        icon_para = icon_frame.paragraphs[0]
        icon_para.font.size = Pt(60)
        icon_para.alignment = PP_ALIGN.CENTER

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(card['x'] + 0.2), Inches(3),
            Inches(2.4), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = card['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['dark']
        title_para.alignment = PP_ALIGN.CENTER

        # 설명
        desc_box = slide.shapes.add_textbox(
            Inches(card['x'] + 0.2), Inches(3.8),
            Inches(2.4), Inches(1.8)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = card['desc']
        desc_frame.word_wrap = True
        for para in desc_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = COLORS['text_light']
            para.alignment = PP_ALIGN.CENTER

    # 하단 장식
    add_decorative_shapes(slide)


def create_code_slide(prs):
    """코드 예시 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "[코드 제목]"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 코드 블록 배경 (어두운 배경)
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(5.5)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(40, 44, 52)  # 어두운 회색
    code_bg.line.fill.background()

    # 코드 영역 (사용자가 코드를 입력할 부분)
    code_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5),
        Inches(8.4), Inches(4.8)
    )
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    code_frame.text = "[여기에 코드를 입력하세요]\n\nfunction example() {\n  // 코드 예시\n  return true;\n}"
    for para in code_frame.paragraphs:
        para.font.name = 'Consolas'
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(171, 178, 191)  # 연한 회색 (코드용)

    # 하단 설명
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.9),
        Inches(9), Inches(0.4)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "💡 [코드 설명을 여기에 작성하세요]"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = COLORS['text_light']


def create_database_slide(prs):
    """데이터베이스/아키텍처 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🗄️ 데이터베이스 구조"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 다이어그램 영역 (연한 배경)
    diagram_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.2),
        Inches(8), Inches(5.3)
    )
    diagram_bg.fill.solid()
    diagram_bg.fill.fore_color.rgb = COLORS['light']
    diagram_bg.line.width = Pt(1)
    diagram_bg.line.color.rgb = COLORS['primary']

    # 다이어그램 안내 텍스트
    diagram_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2),
        Inches(7), Inches(4)
    )
    diagram_frame = diagram_box.text_frame
    diagram_frame.word_wrap = True
    diagram_frame.text = "[여기에 ERD 또는 아키텍처 다이어그램을 삽입하세요]\n\n💡 이미지 삽입 방법:\n1. 슬라이드 편집 모드에서\n2. '삽입' → '이미지' 클릭\n3. ERD 이미지 선택"
    diagram_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in diagram_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = COLORS['text_light']
        para.alignment = PP_ALIGN.CENTER

    # 하단 장식
    add_decorative_shapes(slide)


def create_ending_slide(prs):
    """종료 슬라이드 - 감사 인사"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 그라데이션
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()

    # 중앙 원형 장식
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(2.5), Inches(1.5),
        Inches(5), Inches(5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['secondary']
    circle.line.width = Pt(4)
    circle.line.color.rgb = COLORS['white']

    # 메인 메시지
    main_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.2)
    )
    main_frame = main_box.text_frame
    main_frame.text = "발표를 들어주셔서\n감사합니다! 🍞✨"
    main_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in main_frame.paragraphs:
        para.font.size = Pt(44)
        para.font.bold = True
        para.font.color.rgb = COLORS['white']
        para.alignment = PP_ALIGN.CENTER

    # 하단 정보
    footer_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5),
        Inches(8), Inches(1)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "Breadtopia (빵토피아)\nSSAFY 12기 최종 프로젝트\n\n🤖 AI 기반 빵집 추천 플랫폼"
    for para in footer_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = COLORS['white']
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(6)


if __name__ == "__main__":
    output_path = create_presentation()
    print(f"\nTemplate created successfully!")
    print(f"Path: {output_path}")
    print(f"\nHow to use:")
    print(f"  1. Open file in PowerPoint")
    print(f"  2. Replace [title], [content] placeholders with actual content")
    print(f"  3. Insert images via 'Insert' -> 'Image'")
    print(f"  4. Copy slides as needed")
