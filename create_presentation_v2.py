#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Breadtopia 프로젝트 발표용 PPTX 생성기 v2
- assets/images 및 mascot 이미지 활용
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# 빵토피아 테마 컬러 팔레트
COLORS = {
    'primary': RGBColor(255, 216, 155),      # #FFD89B - 부드러운 빵색
    'secondary': RGBColor(255, 168, 87),     # #FFA857 - 따뜻한 오렌지
    'accent': RGBColor(245, 187, 133),       # #F5BB85 - 크루아상색
    'dark': RGBColor(139, 90, 43),           # #8B5A2B - 갈색 (텍스트)
    'light': RGBColor(255, 248, 238),        # #FFF8EE - 연한 크림색
    'white': RGBColor(255, 255, 255),        # #FFFFFF - 흰색
    'text_dark': RGBColor(51, 51, 51),       # #333333 - 진한 회색
    'text_light': RGBColor(102, 102, 102),   # #666666 - 회색
}

# 이미지 경로
ASSETS_BASE = r"c:\Users\SSAFY\Desktop\최종\final-pjt\frontend\breadmap-project\src\assets"
IMAGES = {
    'logo': os.path.join(ASSETS_BASE, 'images', 'logo.png'),
    'main_screen': os.path.join(ASSETS_BASE, 'images', '메인화면.png'),
    'baguette_corgi': os.path.join(ASSETS_BASE, 'images', '바게트코기.png'),
    'bread': os.path.join(ASSETS_BASE, 'images', '식빵.png'),
    'bread_mouse': os.path.join(ASSETS_BASE, 'images', '식빵쥐.png'),
    'shelf': os.path.join(ASSETS_BASE, 'images', '진열대.png'),
    'shelf2': os.path.join(ASSETS_BASE, 'images', '진열대2.png'),
    'mascot1': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (1).png'),
    'mascot2': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (2).png'),
    'mascot3': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (3).png'),
    'mascot4': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (4).png'),
    'mascot5': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (5).png'),
    'mascot6': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (6).png'),
    'mascot7': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (7).png'),
    'mascot8': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (8).png'),
    'mascot9': os.path.join(ASSETS_BASE, 'mascot', '캐릭터 (9).png'),
}

def create_presentation():
    """빵토피아 발표용 PPTX 생성 (이미지 포함)"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 타이틀 슬라이드 (로고 + 마스코트)
    create_title_slide(prs)

    # 2. 목차 슬라이드 (마스코트 캐릭터)
    create_table_of_contents_slide(prs)

    # 3. 섹션 구분 슬라이드
    create_section_slide(prs, "1. 서비스 소개")

    # 4. 서비스 소개 슬라이드 (메인 화면 이미지)
    create_service_intro_slide(prs)

    # 5. AI 기능 3가지 슬라이드 (마스코트)
    create_ai_features_slide(prs)

    # 6. 2열 레이아웃 (코기 이미지)
    create_two_column_slide(prs)

    # 7. 코드 슬라이드
    create_code_slide(prs)

    # 8. 데이터베이스 구조 슬라이드
    create_database_slide(prs)

    # 9. 종료 슬라이드 (마스코트들)
    create_ending_slide(prs)

    # 저장
    output_path = r"c:\Users\SSAFY\Desktop\1st_PJT_recordings\Breadtopia_Presentation_V2.pptx"
    prs.save(output_path)
    print(f"PPTX file created: {output_path}")
    return output_path


def add_image_if_exists(slide, image_key, left, top, width, height=None):
    """이미지가 존재하면 슬라이드에 추가"""
    if image_key in IMAGES and os.path.exists(IMAGES[image_key]):
        try:
            if height:
                slide.shapes.add_picture(
                    IMAGES[image_key],
                    Inches(left), Inches(top),
                    width=Inches(width), height=Inches(height)
                )
            else:
                slide.shapes.add_picture(
                    IMAGES[image_key],
                    Inches(left), Inches(top),
                    width=Inches(width)
                )
            return True
        except Exception as e:
            print(f"Failed to add image {image_key}: {e}")
            return False
    return False


def create_title_slide(prs):
    """타이틀 슬라이드 - 로고 + 마스코트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['light']
    bg_shape.line.fill.background()

    # 상단 장식 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(2)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()

    # 로고 (중앙 상단)
    if add_image_if_exists(slide, 'logo', 4, 0.3, 2):
        print("Logo added successfully")

    # 프로젝트명
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Breadtopia (빵토피아)"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(56)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.6),
        Inches(8), Inches(0.7)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "빵지순례의 새로운 시작, 당신만의 빵 여행을 시작하세요"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(22)
    subtitle_para.font.color.rgb = COLORS['text_light']
    subtitle_para.alignment = PP_ALIGN.CENTER

    # AI 기능 강조
    ai_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(4.5),
        Inches(5), Inches(0.5)
    )
    ai_frame = ai_box.text_frame
    ai_frame.text = "🤖 AI 기반 3대 서비스 탑재"
    ai_para = ai_frame.paragraphs[0]
    ai_para.font.size = Pt(20)
    ai_para.font.color.rgb = COLORS['secondary']
    ai_para.font.bold = True
    ai_para.alignment = PP_ALIGN.CENTER

    # 좌측 하단 마스코트
    add_image_if_exists(slide, 'bread_mouse', 0.5, 5.5, 1.5)

    # 우측 하단 코기
    add_image_if_exists(slide, 'baguette_corgi', 8, 5.5, 1.5)

    # 하단 정보
    footer_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.8),
        Inches(8), Inches(0.5)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "SSAFY 12기 최종 프로젝트 | 2025.12.24"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(16)
    footer_para.font.color.rgb = COLORS['text_light']
    footer_para.alignment = PP_ALIGN.CENTER


def create_table_of_contents_slide(prs):
    """목차 슬라이드 - 마스코트 캐릭터 배치"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 헤더 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "📌 목차"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 목차 내용
    content_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.5),
        Inches(6.5), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    toc_items = [
        "1. 서비스 소개",
        "2. 개발 배경",
        "3. 기술 스택",
        "4. 시스템 아키텍처 & ERD",
        "5. 핵심 기능 소개 (AI 3대 기능)",
        "6. 필수 요구사항 충족",
        "7. AI 기능 및 추천 알고리즘 상세",
        "8. 게이미피케이션 시스템",
        "9. 주요 코드 설명",
        "10. 프로젝트 회고",
    ]

    for i, item in enumerate(toc_items):
        p = content_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(10)
        p.level = 0

    # 우측 마스코트들
    add_image_if_exists(slide, 'mascot1', 8.2, 1.5, 1.3)
    add_image_if_exists(slide, 'mascot2', 8.2, 3.2, 1.3)
    add_image_if_exists(slide, 'mascot3', 8.2, 4.9, 1.3)


def create_section_slide(prs, section_title):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()

    # 중앙 원형
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3), Inches(2),
        Inches(4), Inches(4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['primary']
    circle.line.width = Pt(3)
    circle.line.color.rgb = COLORS['white']

    # 섹션 제목
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER

    # 좌측 하단 마스코트
    add_image_if_exists(slide, 'mascot4', 0.5, 6, 1.2)

    # 우측 하단 마스코트
    add_image_if_exists(slide, 'mascot5', 8.5, 6, 1.2)


def create_service_intro_slide(prs):
    """서비스 소개 슬라이드 - 메인 화면 이미지 활용"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['light']
    header.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "서비스 핵심 컨셉"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 좌측: 설명
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(4.5), Inches(5.8)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    concepts = [
        "🐶 귀여운 UI/UX",
        "  바구니 속 강아지 빵 컨셉",
        "",
        "🗺️ 지도 기반 빵집 탐색",
        "  Kakao Map API 활용",
        "",
        "🤖 AI 기반 3대 서비스",
        "  • 오늘의 빵집 추천",
        "  • 웰시코기 챗봇",
        "  • 리뷰 요약",
        "",
        "🎮 게이미피케이션",
        "  레벨, 뱃지, 빵 여권",
        "",
        "🤝 소셜 커뮤니티",
        "  빵덕후들의 정보 공유",
    ]

    for concept in concepts:
        p = left_frame.add_paragraph()
        p.text = concept
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(4)

    # 우측: 메인 화면 이미지
    if not add_image_if_exists(slide, 'main_screen', 5.3, 1.2, 4.2):
        # 이미지가 없을 경우 대체 박스
        right_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.3), Inches(1.2),
            Inches(4.2), Inches(5.8)
        )
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = COLORS['light']
        right_bg.line.fill.background()

        placeholder = slide.shapes.add_textbox(
            Inches(5.5), Inches(3.5),
            Inches(3.8), Inches(1.5)
        )
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = "[메인 화면\n스크린샷]"
        placeholder_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for para in placeholder_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = COLORS['text_light']
            para.alignment = PP_ALIGN.CENTER


def create_ai_features_slide(prs):
    """AI 3대 기능 소개 - 마스코트 활용"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🤖 AI 기반 3대 핵심 기능"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    # 3개 카드
    cards = [
        {
            'title': '1. 오늘의 빵집 추천',
            'icon': '🍞',
            'desc': '취향 키워드 기반\n하루 1개 빵집 추천',
            'x': 0.5,
            'mascot': 'mascot6'
        },
        {
            'title': '2. 웰시코기 챗봇',
            'icon': '🐶',
            'desc': '우측 하단 플로팅\n실시간 문의 응대',
            'x': 3.5,
            'mascot': 'mascot7'
        },
        {
            'title': '3. 리뷰 요약',
            'icon': '📝',
            'desc': 'AI가 수십 개 리뷰를\n3-5줄로 요약',
            'x': 6.5,
            'mascot': 'mascot8'
        }
    ]

    for card in cards:
        # 카드 배경
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(card['x']), Inches(1.5),
            Inches(2.8), Inches(5)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLORS['light']
        card_bg.line.width = Pt(2)
        card_bg.line.color.rgb = COLORS['primary']

        # 마스코트 이미지
        if not add_image_if_exists(slide, card['mascot'], card['x'] + 0.6, 1.8, 1.6):
            # 대체 아이콘
            icon_box = slide.shapes.add_textbox(
                Inches(card['x'] + 0.3), Inches(2),
                Inches(2.2), Inches(0.8)
            )
            icon_frame = icon_box.text_frame
            icon_frame.text = card['icon']
            icon_para = icon_frame.paragraphs[0]
            icon_para.font.size = Pt(60)
            icon_para.alignment = PP_ALIGN.CENTER

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(card['x'] + 0.2), Inches(3.7),
            Inches(2.4), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = card['title']
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = COLORS['dark']
        title_para.alignment = PP_ALIGN.CENTER

        # 설명
        desc_box = slide.shapes.add_textbox(
            Inches(card['x'] + 0.2), Inches(4.5),
            Inches(2.4), Inches(1.5)
        )
        desc_frame = desc_box.text_frame
        desc_frame.text = card['desc']
        desc_frame.word_wrap = True
        for para in desc_frame.paragraphs:
            para.font.size = Pt(13)
            para.font.color.rgb = COLORS['text_light']
            para.alignment = PP_ALIGN.CENTER


def create_two_column_slide(prs):
    """2열 레이아웃 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['light']
    header.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "[제목을 입력하세요]"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 좌측 컬럼
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(4.5), Inches(5.8)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.text = "[여기에 내용을 입력하세요]\n\n• 항목 1\n• 항목 2\n• 항목 3"
    for para in left_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = COLORS['text_dark']
        para.space_before = Pt(6)

    # 우측 배경
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.2),
        Inches(4.3), Inches(5.8)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = COLORS['light']
    right_bg.line.fill.background()

    # 우측 이미지 (진열대 또는 플레이스홀더)
    if not add_image_if_exists(slide, 'shelf2', 5.5, 2, 3.7):
        # 플레이스홀더
        placeholder = slide.shapes.add_textbox(
            Inches(5.5), Inches(3),
            Inches(3.8), Inches(3)
        )
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = "[이미지 또는\n다이어그램\n영역]"
        placeholder_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for para in placeholder_frame.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = COLORS['text_light']
            para.alignment = PP_ALIGN.CENTER

    # 하단 마스코트
    add_image_if_exists(slide, 'bread', 8.5, 6.3, 1)


def create_code_slide(prs):
    """코드 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "[코드 제목 - 예: AI 추천 알고리즘]"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 코드 블록 배경
    code_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(5.5)
    )
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(40, 44, 52)
    code_bg.line.fill.background()

    # 코드 영역
    code_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5),
        Inches(8.4), Inches(4.8)
    )
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    code_frame.text = "[여기에 코드를 입력하세요]\n\nfunction getTodayRecommendedStore() {\n  // 코드 예시\n  return recommendedStore;\n}"
    for para in code_frame.paragraphs:
        para.font.name = 'Consolas'
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(171, 178, 191)

    # 하단 설명
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.9),
        Inches(9), Inches(0.4)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "💡 [코드 설명을 여기에 작성하세요]"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = COLORS['text_light']


def create_database_slide(prs):
    """데이터베이스 구조 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🗄️ 데이터베이스 구조 & 공공데이터 활용"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(30)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    # 좌측: ERD 영역
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.2),
        Inches(5.5), Inches(5.5)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = COLORS['light']
    left_bg.line.width = Pt(1)
    left_bg.line.color.rgb = COLORS['primary']

    left_text = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(4.5), Inches(3)
    )
    left_frame = left_text.text_frame
    left_frame.word_wrap = True
    left_frame.text = "[ERD 다이어그램\n삽입 영역]\n\n💡 User - Review - Store\n관계 표시"
    left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in left_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = COLORS['text_light']
        para.alignment = PP_ALIGN.CENTER

    # 우측: 데이터 수집 전략
    right_box = slide.shapes.add_textbox(
        Inches(6.2), Inches(1.2),
        Inches(3.3), Inches(5.5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    data_points = [
        "📊 데이터 수집 전략",
        "",
        "1️⃣ 공공데이터포털",
        "   전국 제과점 정보",
        "",
        "2️⃣ 경위도 필수 모델",
        "   latitude, longitude",
        "",
        "3️⃣ Geocoding 활용",
        "   주소 → 좌표 변환",
        "   성공률 95%",
        "",
        "4️⃣ 데이터 규모",
        "   빵집 500개+",
        "   리뷰 500개+",
    ]

    for point in data_points:
        p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(3)

    # 하단 마스코트
    add_image_if_exists(slide, 'mascot9', 0.3, 6.3, 1)


def create_ending_slide(prs):
    """종료 슬라이드 - 마스코트들과 함께"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()

    # 중앙 원형
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(2.5), Inches(1.5),
        Inches(5), Inches(5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['secondary']
    circle.line.width = Pt(4)
    circle.line.color.rgb = COLORS['white']

    # 메인 메시지
    main_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.2)
    )
    main_frame = main_box.text_frame
    main_frame.text = "발표를 들어주셔서\n감사합니다!"
    main_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in main_frame.paragraphs:
        para.font.size = Pt(44)
        para.font.bold = True
        para.font.color.rgb = COLORS['white']
        para.alignment = PP_ALIGN.CENTER

    # 하단 정보
    footer_box = slide.shapes.add_textbox(
        Inches(1), Inches(5),
        Inches(8), Inches(1.2)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "Breadtopia (빵토피아)\nSSAFY 12기 최종 프로젝트\n\n🤖 AI 기반 빵집 추천 플랫폼"
    for para in footer_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = COLORS['white']
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(6)

    # 양쪽 하단 마스코트
    add_image_if_exists(slide, 'bread_mouse', 0.3, 6.3, 1.2)
    add_image_if_exists(slide, 'baguette_corgi', 8.7, 6.3, 1.2)

    # 상단 양쪽 마스코트
    add_image_if_exists(slide, 'mascot1', 0.3, 0.3, 1)
    add_image_if_exists(slide, 'mascot2', 8.9, 0.3, 1)


if __name__ == "__main__":
    print("Creating presentation with images...")
    output_path = create_presentation()
    print(f"\nTemplate created successfully!")
    print(f"Path: {output_path}")
    print(f"\nFeatures:")
    print(f"  - Logo and mascot characters included")
    print(f"  - Breadtopia theme colors")
    print(f"  - AI features slide with mascots")
    print(f"  - Main screen preview")
    print(f"\nHow to use:")
    print(f"  1. Open file in PowerPoint")
    print(f"  2. Replace placeholders with actual content")
    print(f"  3. Add more slides by copying templates")
    print(f"  4. Insert screenshots where needed")
